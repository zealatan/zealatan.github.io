"""
Generate slides/project_overview.pptx using ONLY native python-pptx elements.
No matplotlib images — every shape, text box, connector, chart, and table
is a first-class PowerPoint object that can be selected and edited.
"""

import os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn

# ── palette ───────────────────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x1A, 0x2E, 0x4A)
C_TEAL   = RGBColor(0x00, 0x8B, 0x8B)
C_ORANGE = RGBColor(0xE8, 0x7D, 0x1E)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xF2, 0xF4, 0xF6)
C_MGRAY  = RGBColor(0xCC, 0xCC, 0xCC)
C_DGRAY  = RGBColor(0x44, 0x44, 0x44)
C_GREEN  = RGBColor(0x27, 0xAE, 0x60)
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_BLUE   = RGBColor(0x2C, 0x5F, 0x8A)
C_LBLUE  = RGBColor(0x1A, 0x6E, 0x8B)
C_STEEL  = RGBColor(0x2C, 0x3E, 0x50)
FONT     = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]   # completely blank

# ── primitive helpers ─────────────────────────────────────────────────────────

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill, lc=None, lw=0.5, rounded=False):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as A
    stype = A.ROUNDED_RECTANGLE if rounded else A.RECTANGLE
    s = slide.shapes.add_shape(stype, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if lc:
        s.line.color.rgb = lc
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def add_txt(slide, text, x, y, w, h,
            size=12, bold=False, italic=False,
            color=C_DGRAY, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name   = FONT
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def set_shape_text(shape, lines, size=11, color=C_WHITE,
                   align=PP_ALIGN.CENTER, bold_first=True):
    """Write one or more lines into a shape's text frame (vertically centred)."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf._txBody.set('anchor', 'ctr')
    tf.margin_left   = Inches(0.06)
    tf.margin_right  = Inches(0.06)
    tf.margin_top    = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    items = lines if isinstance(lines, list) else [lines]
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(1)
        r = p.add_run()
        r.text = line
        r.font.name  = FONT
        r.font.size  = Pt(size)
        r.font.bold  = bold_first and (i == 0)
        r.font.color.rgb = color


def add_box(slide, x, y, w, h, lines, fill,
            tc=C_WHITE, size=11, bf=True, rounded=True, lc=None):
    s = add_rect(slide, x, y, w, h, fill,
                 lc=lc or RGBColor(0x99, 0x99, 0x99), lw=0.75, rounded=rounded)
    set_shape_text(s, lines, size=size, color=tc, bold_first=bf)
    return s


def add_arrow(slide, x1, y1, x2, y2, color=C_ORANGE, width=1.5):
    """Straight connector with an arrowhead at (x2, y2)."""
    c = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    # Inject tailEnd arrowhead into the XML
    spPr = c._element.find(qn('p:spPr'))
    if spPr is not None:
        ln = spPr.find(qn('a:ln'))
        if ln is not None:
            te = etree.SubElement(ln, qn('a:tailEnd'))
            te.set('type', 'arrow')
            te.set('w', 'med')
            te.set('len', 'med')
    return c


def add_header(slide, title, size=28):
    add_rect(slide, 0, 0, 13.33, 1.05, fill=C_NAVY)
    add_rect(slide, 0, 1.00, 13.33, 0.07, fill=C_ORANGE)
    add_txt(slide, title, 0.35, 0.12, 12.6, 0.82,
            size=size, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)


def add_stat_card(slide, x, y, value, label, fill=C_TEAL):
    add_rect(slide, x, y, 2.5, 1.4, fill=fill, rounded=True,
             lc=RGBColor(0xFF, 0xFF, 0xFF), lw=0)
    add_txt(slide, value, x, y + 0.08, 2.5, 0.72,
            size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_txt(slide, label, x, y + 0.85, 2.5, 0.45,
            size=11, color=C_LGRAY, align=PP_ALIGN.CENTER)


# ── chart helpers ─────────────────────────────────────────────────────────────

def _color_points(series, colors):
    """Set individual data-point fill colors by injecting c:dPt XML nodes."""
    ser = series._element
    cat = ser.find(qn('c:cat'))
    val = ser.find(qn('c:val'))
    ref = cat if cat is not None else val
    pos = list(ser).index(ref) if ref is not None else len(list(ser))
    for i, rgb in enumerate(colors):
        dPt = etree.Element(qn('c:dPt'))
        etree.SubElement(dPt, qn('c:idx')).set('val', str(i))
        etree.SubElement(dPt, qn('c:invertIfNegative')).set('val', '0')
        spPr = etree.SubElement(dPt, qn('c:spPr'))
        sf   = etree.SubElement(spPr, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr')).set('val', str(rgb))
        etree.SubElement(etree.SubElement(spPr, qn('a:ln')), qn('a:noFill'))
        ser.insert(pos + i, dPt)


def add_bar_chart(slide, x, y, w, h, cats, vals, colors,
                  title='', horizontal=True):
    cd = ChartData()
    cd.categories = cats
    cd.add_series('', vals)
    ct = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    cf = slide.shapes.add_chart(ct, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = cf.chart
    ch.has_legend = False
    if title:
        ch.has_title = True
        ch.chart_title.text_frame.text = title
        ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
        ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
        ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = FONT
    _color_points(ch.series[0], colors)
    try:
        ch.category_axis.tick_labels.font.size  = Pt(8)
        ch.category_axis.tick_labels.font.name  = FONT
        ch.value_axis.tick_labels.font.size     = Pt(8)
        ch.value_axis.tick_labels.font.name     = FONT
        ch.plots[0].has_data_labels             = True
        ch.plots[0].data_labels.show_value      = True
        ch.plots[0].data_labels.font.size       = Pt(8)
        ch.plots[0].data_labels.font.name       = FONT
    except Exception:
        pass
    return cf


def add_pie_chart(slide, x, y, w, h, cats, vals, colors, title=''):
    cd = ChartData()
    cd.categories = cats
    cd.add_series('', vals)
    cf = slide.shapes.add_chart(XL_CHART_TYPE.PIE,
                                 Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = cf.chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    if title:
        ch.has_title = True
        ch.chart_title.text_frame.text = title
        ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
        ch.chart_title.text_frame.paragraphs[0].runs[0].font.bold = True
        ch.chart_title.text_frame.paragraphs[0].runs[0].font.name = FONT
    _color_points(ch.series[0], colors)
    try:
        ch.plots[0].has_data_labels            = True
        ch.plots[0].data_labels.show_percentage = True
        ch.plots[0].data_labels.show_category_name = True
        ch.plots[0].data_labels.font.size       = Pt(8)
        ch.plots[0].data_labels.font.name       = FONT
    except Exception:
        pass
    return cf


def add_table(slide, x, y, w, h, headers, rows,
              hdr_fill=C_NAVY, hdr_tc=C_WHITE,
              alt_fill=RGBColor(0xE8, 0xF4, 0xF8)):
    ncols, nrows = len(headers), len(rows) + 1
    tf = slide.shapes.add_table(
        nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h))
    t = tf.table
    # header
    for j, hdr in enumerate(headers):
        c = t.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = hdr_fill
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = hdr
        r.font.name = FONT; r.font.size = Pt(9)
        r.font.bold = True; r.font.color.rgb = hdr_tc
    # data rows
    for i, row in enumerate(rows):
        bg = alt_fill if i % 2 == 0 else C_WHITE
        for j, val in enumerate(row):
            c = t.cell(i + 1, j)
            c.fill.solid(); c.fill.fore_color.rgb = bg
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(val)
            r.font.name = FONT; r.font.size = Pt(8.5)
            if val == 'PASS':
                r.font.color.rgb = C_GREEN; r.font.bold = True
            elif val == 'FAIL':
                r.font.color.rgb = C_RED; r.font.bold = True
            else:
                r.font.color.rgb = C_DGRAY
    return tf


# ── waveform helper ───────────────────────────────────────────────────────────

def add_waveform(slide, x0, y0, w, h, signals, state_labels=None,
                 label_w=1.1, max_font=8.5):
    """
    Draw a native-pptx digital timing diagram (no images).
    signals    : [(name, [0/1 per cycle], RGBColor or None)]
    state_labels: optional list of FSM-state strings, one per cycle
    """
    HDR_H   = 0.30 if state_labels else 0.22
    n_sig   = len(signals)
    n_cyc   = len(signals[0][1])
    sig_x0  = x0 + label_w
    sig_w   = w  - label_w
    row_h   = (h - HDR_H) / n_sig
    cyc_w   = sig_w / n_cyc

    DARK_BG  = RGBColor(0x0D, 0x1B, 0x2A)
    LBL_BG   = RGBColor(0x0A, 0x14, 0x1E)
    EVEN_COL = RGBColor(0x12, 0x22, 0x38)
    GRID_CLR = RGBColor(0x22, 0x33, 0x4A)
    NUM_CLR  = RGBColor(0x77, 0x99, 0xBB)

    # Background
    add_rect(slide, x0, y0, w, h, fill=DARK_BG, lc=GRID_CLR, lw=0.75)
    add_rect(slide, x0, y0, label_w, h, fill=LBL_BG)

    # Alternating column shading + cycle numbers
    for c in range(n_cyc):
        cx = sig_x0 + c * cyc_w
        add_rect(slide, cx, y0, cyc_w, h,
                 fill=EVEN_COL if c % 2 == 0 else DARK_BG)
        add_txt(slide, str(c + 1),
                cx, y0 + 0.02, cyc_w, 0.17,
                size=8, bold=True, color=NUM_CLR, align=PP_ALIGN.CENTER)
        if state_labels and c < len(state_labels):
            add_txt(slide, state_labels[c],
                    cx + 0.02, y0 + 0.17, cyc_w - 0.04, 0.12,
                    size=6, color=C_ORANGE, align=PP_ALIGN.CENTER)

    # Corner labels
    add_txt(slide, "cyc →",
            x0 + 0.04, y0 + 0.02, label_w - 0.08, 0.17,
            size=7, italic=True, color=NUM_CLR, align=PP_ALIGN.RIGHT)
    if state_labels:
        add_txt(slide, "state:",
                x0 + 0.04, y0 + 0.17, label_w - 0.08, 0.12,
                size=6, italic=True, color=C_ORANGE, align=PP_ALIGN.RIGHT)

    # Header divider (orange accent)
    add_rect(slide, x0, y0 + HDR_H, w, 0.010, fill=C_ORANGE)

    # Vertical grid lines at cycle edges
    for c in range(1, n_cyc):
        gx = sig_x0 + c * cyc_w
        add_rect(slide, gx - 0.004, y0, 0.008, h, fill=GRID_CLR)

    # Each signal row
    for si, (name, values, color) in enumerate(signals):
        if color is None:
            color = C_TEAL
        row_y  = y0 + HDR_H + si * row_h
        hi_top = row_y + row_h * 0.08
        hi_h   = row_h * 0.62
        lo_y   = row_y + row_h * 0.80
        lo_h   = max(0.020, row_h * 0.08)

        # Row separator
        add_rect(slide, x0, row_y, w, 0.006, fill=GRID_CLR)

        # Signal name label
        add_txt(slide, name,
                x0 + 0.04, row_y + row_h * 0.12,
                label_w - 0.10, row_h * 0.76,
                size=min(max_font, row_h * 26),
                bold=True, color=color, align=PP_ALIGN.RIGHT)

        # Waveform segments
        prev_val = None
        for c, val in enumerate(values):
            cx_left = sig_x0 + c * cyc_w
            seg_x   = cx_left + 0.015
            seg_w_s = cyc_w   - 0.030
            if val:
                add_rect(slide, seg_x, hi_top, seg_w_s, hi_h, fill=color)
            else:
                add_rect(slide, seg_x, lo_y, seg_w_s, lo_h, fill=color)
            if prev_val is not None and val != prev_val:
                span = (lo_y + lo_h) - hi_top
                add_rect(slide, cx_left - 0.010, hi_top, 0.020, span, fill=color)
            prev_val = val


# =============================================================================
# Slide 1 – Title
# =============================================================================
def slide_title():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, C_NAVY)

    # Accent stripe
    add_rect(sl, 0, 3.1, 13.33, 0.08, fill=C_ORANGE)

    # Title
    add_txt(sl, "AXI4 RTL Verification Project",
            0.7, 1.1, 11.9, 1.4,
            size=44, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_txt(sl, "AXI-Lite Register File  ·  AXI4 Memory Model  ·  Simple AXI4 Master",
            0.7, 2.55, 11.9, 0.6,
            size=20, color=C_TEAL, align=PP_ALIGN.CENTER)

    add_txt(sl, "Vivado xsim  ·  SystemVerilog testbenches  ·  Automated CI gate",
            0.7, 3.3, 11.9, 0.5,
            size=15, italic=True,
            color=RGBColor(0x99, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

    # Four stat cards
    for (val, lab, fill), x in zip(
        [("3",   "RTL / Model\nComponents",   C_TEAL),
         ("177", "Total Checks\nPassed",       C_GREEN),
         ("0",   "Failures",                   C_BLUE),
         ("3",   "CI Gate\nScripts",           C_LBLUE)],
        [1.2, 3.95, 6.7, 9.45]):
        add_stat_card(sl, x, 4.0, val, lab, fill=fill)

    add_txt(sl, "Bumhee Lee  ·  2026-05-01",
            0, 6.9, 13.33, 0.4,
            size=12, italic=True,
            color=RGBColor(0x77, 0x88, 0x99), align=PP_ALIGN.CENTER)


slide_title()


# =============================================================================
# Slide 2 – Project Architecture (block diagram, all native shapes)
# =============================================================================
def slide_arch():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, C_LGRAY)
    add_header(sl, "Project Architecture Overview")

    # ── Three column headers ─────────────────────────────────────────
    COL = [0.35, 4.55, 8.75]   # left edge of each column
    CW  = 3.9                  # column width

    for (title, sub), cx in zip(
        [("AXI-Lite Register File", "rtl/axi_lite_regfile.v"),
         ("AXI4 Memory Model",      "tb/axi_mem_model.sv"),
         ("Simple AXI4 Master",     "rtl/simple_axi_master.v")],
        COL):
        add_rect(sl, cx, 1.15, CW, 0.55, fill=C_NAVY,
                 lc=C_TEAL, lw=1.0, rounded=True)
        add_txt(sl, title, cx, 1.17, CW, 0.32,
                size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_txt(sl, sub,   cx, 1.46, CW, 0.22,
                size=9, italic=True,
                color=RGBColor(0x99, 0xCC, 0xDD), align=PP_ALIGN.CENTER)

    # ── Testbench boxes (row 1) ──────────────────────────────────────
    TB_Y = 1.95
    for (label, checks), cx in zip(
        [("axi_lite_regfile_tb.sv\n133 checks", "133 / 133 PASS"),
         ("mem_rw_tb.sv\n30 checks",            "30 / 30 PASS"),
         ("simple_axi_master_tb.sv\n14 checks", "14 / 14 PASS")],
        COL):
        add_box(sl, cx, TB_Y, CW, 0.85,
                [label.split('\n')[0], label.split('\n')[1]],
                fill=C_TEAL, tc=C_WHITE, size=10)

    # ── RTL / Model boxes (row 2) ────────────────────────────────────
    RTL_Y = 3.3
    for (lines, fill), cx in zip(
        [
            (["4 × 32-bit regs",
              "SLVERR on OOB addr",
              "3-path write FSM"], C_BLUE),
            (["1 KB RAM, Parameterised",
              "Byte-WSTRB support",
              "SLVERR on OOB addr"], C_LBLUE),
            (["6-state FSM",
              "Write → Readback",
              "Error flag detection"], C_STEEL),
        ],
        COL):
        add_box(sl, cx, RTL_Y, CW, 1.2, lines, fill=fill, size=10)

    # ── Vertical arrows: TB → RTL ────────────────────────────────────
    for cx in COL:
        mid = cx + CW / 2
        add_arrow(sl, mid, TB_Y + 0.85, mid, RTL_Y, color=C_ORANGE, width=1.5)

    # ── Horizontal arrow: Master ↔ Memory (AXI4 bus) ─────────────────
    # From master's left edge to memory model's right edge
    add_arrow(sl, COL[2],       RTL_Y + 0.6,
              COL[1] + CW + 0.02, RTL_Y + 0.6,
              color=C_TEAL, width=2.0)
    add_txt(sl, "AXI4 bus", COL[1] + CW + 0.06, RTL_Y + 0.35, 0.95, 0.3,
            size=9, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)

    # ── CI Script boxes (row 3) ──────────────────────────────────────
    CI_Y = 5.0
    for (script,), cx in zip(
        [("run_axi_regfile_sim.sh",),
         ("run_mem_rw_sim.sh",),
         ("run_simple_axi_master_sim.sh",)],
        COL):
        add_box(sl, cx, CI_Y, CW, 0.65, [script], fill=C_NAVY, size=9)

    # ── Vertical arrows: RTL → CI ─────────────────────────────────────
    for cx in COL:
        mid = cx + CW / 2
        add_arrow(sl, mid, RTL_Y + 1.2, mid, CI_Y, color=C_MGRAY, width=1.2)

    # ── CI gate legend ────────────────────────────────────────────────
    add_rect(sl, 0.35, 5.85, 12.63, 0.6, fill=C_NAVY, rounded=True)
    add_txt(sl, "CI gate:  grep -qE '\\[FAIL\\]|FATAL' <log>  →  exit 1 on any failure",
            0.5, 5.93, 12.3, 0.42,
            size=12, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)


slide_arch()


# =============================================================================
# Slide 3 – AXI-Lite Register File: block diagram + write FSM
# =============================================================================
def slide_regfile():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, C_LGRAY)
    add_header(sl, "AXI-Lite Register File  (rtl/axi_lite_regfile.v)")

    # ── LEFT: Block Diagram ──────────────────────────────────────────
    add_txt(sl, "Block Diagram", 0.3, 1.15, 5.5, 0.35,
            size=13, bold=True, color=C_NAVY, align=PP_ALIGN.LEFT)

    # AXI channel boxes
    CH_X = 0.3
    for label, y, fill in [
        ("AW channel\nawaddr / awvalid / awready", 1.6,  C_TEAL),
        ("W channel\nwdata / wstrb / wvalid / wready", 2.55, C_TEAL),
        ("B channel\nbresp / bvalid / bready",     3.5,  C_LBLUE),
        ("AR channel\naraddr / arvalid / arready",  4.45, C_TEAL),
        ("R channel\nrdata / rresp / rvalid / rready", 5.4, C_LBLUE),
    ]:
        add_box(sl, CH_X, y, 2.0, 0.75, label.split('\n'), fill=fill, size=8.5)

    # DUT box
    add_box(sl, 2.55, 2.4, 1.6, 2.2,
            ["axi_lite\n_regfile", "", "regs[0..3]", "32-bit × 4"],
            fill=C_NAVY, size=9)

    # Register boxes
    for i in range(4):
        add_box(sl, 4.4, 1.6 + i * 0.95, 1.3, 0.75,
                [f"reg[{i}]", f"0x{i*4:02X}"],
                fill=C_BLUE, size=9)

    # Arrows: channels → DUT
    for y_mid in [1.975, 2.925]:
        add_arrow(sl, 2.3, y_mid, 2.55, y_mid, color=C_ORANGE)
    add_arrow(sl, 2.55, 3.5, 2.3, 3.5, color=C_TEAL)   # B channel (response)
    for y_mid in [4.825, 5.775]:
        add_arrow(sl, 2.3, y_mid, 2.55, y_mid, color=C_ORANGE)
    add_arrow(sl, 2.55, 5.0, 2.3, 5.0, color=C_TEAL)   # R channel (response)

    # Arrows: DUT → registers
    for i in range(4):
        add_arrow(sl, 4.15, 1.975 + i * 0.95, 4.4, 1.975 + i * 0.95,
                  color=C_ORANGE)

    # SLVERR note
    add_rect(sl, 0.3, 6.2, 5.5, 0.55, fill=C_RED, rounded=True)
    add_txt(sl, "addr[31:4] ≠ 0  →  SLVERR, register unchanged",
            0.35, 6.27, 5.4, 0.38,
            size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ── RIGHT: Write FSM ─────────────────────────────────────────────
    add_txt(sl, "Write FSM  (4 states)", 6.5, 1.15, 6.5, 0.35,
            size=13, bold=True, color=C_NAVY, align=PP_ALIGN.LEFT)

    # State boxes arranged in a 2×2 grid
    #   W_IDLE   (top-left)   W_WAIT_W (top-right)
    #   W_WAIT_A (bot-left)   W_BRESP  (bot-right)
    FSM = {
        "W_IDLE":   (6.6,  2.1),
        "W_WAIT_W": (9.5,  2.1),
        "W_WAIT_A": (6.6,  4.2),
        "W_BRESP":  (9.5,  4.2),
    }
    COLORS = {
        "W_IDLE":   C_STEEL,
        "W_WAIT_W": C_BLUE,
        "W_WAIT_A": C_BLUE,
        "W_BRESP":  C_TEAL,
    }
    DESC = {
        "W_IDLE":   ["W_IDLE", "ready for\nnew txn"],
        "W_WAIT_W": ["W_WAIT_W", "have AW,\nwait W"],
        "W_WAIT_A": ["W_WAIT_A", "have W,\nwait AW"],
        "W_BRESP":  ["W_BRESP", "sending\nresponse"],
    }
    BW, BH = 2.0, 0.95
    for name, (sx, sy) in FSM.items():
        add_box(sl, sx, sy, BW, BH, DESC[name], fill=COLORS[name], size=10)

    cx = {n: x + BW/2 for n, (x, _) in FSM.items()}
    cy = {n: y + BH/2 for n, (_, y) in FSM.items()}

    # W_IDLE → W_WAIT_W (right)
    add_arrow(sl, FSM["W_IDLE"][0]+BW, cy["W_IDLE"],
              FSM["W_WAIT_W"][0], cy["W_WAIT_W"], color=C_ORANGE)
    add_txt(sl, "awvalid only", cx["W_IDLE"]+0.85, cy["W_IDLE"]-0.38, 1.3, 0.3,
            size=8, color=C_DGRAY, align=PP_ALIGN.CENTER)

    # W_IDLE → W_WAIT_A (down)
    add_arrow(sl, cx["W_IDLE"], FSM["W_IDLE"][1]+BH,
              cx["W_WAIT_A"], FSM["W_WAIT_A"][1], color=C_ORANGE)
    add_txt(sl, "wvalid only", FSM["W_IDLE"][0]-1.35, cy["W_IDLE"]+0.95, 1.3, 0.3,
            size=8, color=C_DGRAY, align=PP_ALIGN.CENTER)

    # W_WAIT_W → W_BRESP (down)
    add_arrow(sl, cx["W_WAIT_W"], FSM["W_WAIT_W"][1]+BH,
              cx["W_BRESP"], FSM["W_BRESP"][1], color=C_ORANGE)
    add_txt(sl, "wvalid", cx["W_WAIT_W"]+0.12, cy["W_WAIT_W"]+0.9, 0.9, 0.3,
            size=8, color=C_DGRAY, align=PP_ALIGN.LEFT)

    # W_WAIT_A → W_BRESP (right)
    add_arrow(sl, FSM["W_WAIT_A"][0]+BW, cy["W_WAIT_A"],
              FSM["W_BRESP"][0], cy["W_BRESP"], color=C_ORANGE)
    add_txt(sl, "awvalid", cx["W_WAIT_A"]+0.85, cy["W_WAIT_A"]-0.38, 1.0, 0.3,
            size=8, color=C_DGRAY, align=PP_ALIGN.CENTER)

    # W_IDLE → W_BRESP (both valid, diagonal text)
    add_arrow(sl, FSM["W_IDLE"][0]+BW, FSM["W_IDLE"][1]+BH,
              FSM["W_BRESP"][0], FSM["W_BRESP"][1], color=C_GREEN, width=1.8)
    add_txt(sl, "both valid", cx["W_IDLE"]+1.25, cy["W_IDLE"]+0.7, 1.1, 0.3,
            size=8, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

    # W_BRESP → W_IDLE (back to top-left, via top edge)
    add_arrow(sl, cx["W_BRESP"], FSM["W_BRESP"][1],
              cx["W_IDLE"], FSM["W_IDLE"][1]+BH, color=C_MGRAY, width=1.2)
    add_txt(sl, "bready", cx["W_IDLE"]+2.1, cy["W_IDLE"]+1.15, 0.9, 0.3,
            size=8, color=C_DGRAY, align=PP_ALIGN.CENTER)

    # ── Write-handshake waveform (lower-right, W_IDLE both-valid path) ──
    add_txt(sl, "AXI Write Handshake  (W_IDLE path, awvalid & wvalid simultaneous)",
            6.5, 5.30, 6.33, 0.23,
            size=9, bold=True, color=C_NAVY, align=PP_ALIGN.LEFT)
    # awready/wready pre-asserted (slave in W_IDLE), dip LOW in W_BRESP
    add_waveform(sl, 6.5, 5.53, 6.33, 1.12,
                 [("awvalid", [0, 1, 0, 0, 0, 0], C_TEAL),
                  ("awready", [1, 1, 0, 1, 1, 1], C_GREEN),
                  ("wvalid",  [0, 1, 0, 0, 0, 0], C_BLUE),
                  ("bvalid",  [0, 0, 1, 0, 0, 0], C_LBLUE),
                  ("bready",  [0, 0, 1, 0, 0, 0], C_ORANGE)],
                 state_labels=["·", "W_IDLE", "W_BRESP", "W_IDLE", "·", "·"],
                 label_w=0.90, max_font=7.0)

    # Bottom strip
    add_rect(sl, 0, 6.75, 13.33, 0.75, fill=C_NAVY)
    add_txt(sl, "133 / 133 checks PASSED  ·  "
            "Byte-WSTRB  ·  3 write paths  ·  SLVERR for addr[31:4]≠0  ·  Sim time 3720 ns",
            0.3, 6.82, 12.7, 0.55,
            size=12, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)


slide_regfile()


# =============================================================================
# Slide 4 – AXI-Lite Coverage Charts (native bar + pie)
# =============================================================================
def slide_regfile_charts():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, C_LGRAY)
    add_header(sl, "AXI-Lite Register File — Verification Coverage (133 Checks)")

    cats = ["T1: Reset defaults",   "T2: Write all-ones",
            "T3: Unique values",    "T4: Write all-zeros",
            "T5: Partial WSTRB",    "T5b: Byte isolation",
            "T6: AW-before-W",      "T7: W-before-AW",
            "T8: Invalid addr 0x10","T9: Same-cycle AW+W",
            "T10: B backpressure",  "T11: R backpressure",
            "T12: OOB sweep W_IDLE","T13: OOB AW-first",
            "T14: OOB W-first"]
    vals = (8, 12, 12, 12, 10, 12, 4, 4, 5, 4, 8, 8, 15, 9, 9)
    bar_colors = (
        [C_BLUE]*4 + [C_TEAL]*2 + [C_BLUE]*2 +
        [C_ORANGE] + [C_BLUE]*3 + [C_RED]*3
    )

    # Left: horizontal bar chart
    add_bar_chart(sl, 0.3, 1.15, 6.8, 6.1,
                  cats, vals, bar_colors,
                  title="Checks per Test  (click chart to edit data)",
                  horizontal=True)

    # Right: pie chart by category
    pie_cats = ["Functional (T1–T4, T6–T7, T9)",
                "WSTRB (T5, T5b)",
                "SLVERR (T8, T12–T14)",
                "Backpressure (T10–T11)"]
    pie_vals = (8+12+12+12+4+4+4, 10+12, 5+15+9+9, 8+8)
    add_pie_chart(sl, 7.3, 1.15, 5.7, 5.3,
                  pie_cats, pie_vals,
                  [C_BLUE, C_TEAL, C_RED, C_ORANGE],
                  title="Category Breakdown")

    # Legend note
    for (label, color), y in zip(
        [("Functional tests", C_BLUE), ("WSTRB tests", C_TEAL),
         ("SLVERR tests", C_RED), ("Backpressure tests", C_ORANGE)],
        [6.55, 6.8, 7.05, 7.3]):
        add_rect(sl, 7.3, y - 0.12, 0.25, 0.25, fill=color, rounded=True)
        add_txt(sl, label, 7.65, y - 0.14, 3.5, 0.28, size=10, color=C_DGRAY)


slide_regfile_charts()


# =============================================================================
# Slide 5 – AXI4 Memory Model
# =============================================================================
def slide_mem_model():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, C_LGRAY)
    add_header(sl, "AXI4 Memory Model  (tb/axi_mem_model.sv)")

    # ── Left: block diagram ───────────────────────────────────────────
    add_txt(sl, "Block Diagram", 0.3, 1.15, 6.0, 0.35,
            size=13, bold=True, color=C_NAVY)

    # AXI input channels
    for label, y, fill in [
        ("AW  awaddr / awvalid / awready", 1.65, C_TEAL),
        ("W   wdata / wstrb / wlast / wvalid / wready", 2.5, C_TEAL),
        ("AR  araddr / arvalid / arready", 4.35, C_LBLUE),
    ]:
        add_box(sl, 0.3, y, 2.5, 0.65, [label], fill=fill, size=8.5)

    # AXI output channels
    for label, y, fill in [
        ("B   bresp / bvalid / bready", 3.35, C_TEAL),
        ("R   rdata / rresp / rlast / rvalid / rready", 5.2, C_LBLUE),
    ]:
        add_box(sl, 0.3, y, 2.5, 0.65, [label], fill=fill, size=8.5)

    # Write FSM + Read FSM + memory
    add_box(sl, 3.05, 1.65, 1.4, 1.5,
            ["Write FSM", "4 states", "W_IDLE", "W_WAIT_W/A", "W_BRESP"],
            fill=C_NAVY, size=8.5)
    add_box(sl, 3.05, 3.5, 1.4, 1.5,
            ["Read FSM", "2 states", "R_IDLE", "R_RVALID"],
            fill=C_NAVY, size=8.5)
    add_box(sl, 4.65, 1.65, 1.6, 3.35,
            ["mem[ ]", "1024 × 32-bit", "(4 KB)", "", "init 0x0", "at sim start"],
            fill=C_BLUE, size=9)

    # Arrows: channels → FSMs
    for y_mid in [1.975, 2.825]:
        add_arrow(sl, 2.8, y_mid, 3.05, y_mid, color=C_ORANGE)
    add_arrow(sl, 3.05, 3.0, 2.8, 3.0, color=C_TEAL)          # B output
    add_arrow(sl, 2.8, 4.675, 3.05, 4.675, color=C_ORANGE)    # AR
    add_arrow(sl, 3.05, 5.525, 2.8, 5.525, color=C_TEAL)      # R output

    # Arrows: FSMs ↔ memory
    add_arrow(sl, 4.45, 2.4, 4.65, 2.4, color=C_ORANGE)       # write
    add_arrow(sl, 4.65, 4.3, 4.45, 4.3, color=C_LBLUE)        # read

    # Parameters box
    add_rect(sl, 0.3, 6.2, 6.0, 0.55, fill=C_NAVY, rounded=True)
    add_txt(sl, "Parameters:  MEM_DEPTH (words, default 1024)  ·  MEM_BASE (default 0x0)",
            0.4, 6.27, 5.8, 0.38,
            size=10, color=C_TEAL, align=PP_ALIGN.CENTER)

    # ── Right: test results chart ─────────────────────────────────────
    add_bar_chart(sl, 6.6, 1.15, 6.4, 5.8,
                  ["T1: Full word\nwrite/readback",
                   "T2: Multiple\naddresses",
                   "T3: Byte-lane\nWSTRB",
                   "T4: Unwritten\naddr = 0"],
                  (6, 12, 10, 2),
                  [C_BLUE, C_LBLUE, C_TEAL, C_ORANGE],
                  title="Test Coverage — 30 Checks  (click to edit)",
                  horizontal=False)

    add_rect(sl, 0, 6.8, 13.33, 0.7, fill=C_NAVY)
    add_txt(sl, "30 / 30 checks PASSED  ·  SLVERR for OOB addresses  ·  Simulation time 690 ns",
            0.3, 6.87, 12.7, 0.5,
            size=12, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)


slide_mem_model()


# =============================================================================
# Slide 6 – Simple AXI4 Master: 6-state FSM
# =============================================================================
def slide_master_fsm():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, C_LGRAY)
    add_header(sl, "Simple AXI4 Master — 6-State FSM  (rtl/simple_axi_master.v)")

    # ── FSM states in a 2×3 grid ──────────────────────────────────────
    #   row 0 (y=1.5): IDLE → WR_ADDR → WR_RESP
    #   row 1 (y=3.5): DONE ← RD_DATA ← RD_ADDR
    #
    #   cols: 0.4, 2.5, 4.6   (x of left edge, box width 1.8)

    R0Y, R1Y = 1.5, 3.5    # row top edges
    CX = [0.4, 2.5, 4.6]   # column left edges
    BW, BH = 1.8, 0.9

    # State definitions
    STATES = {
        "IDLE":    (CX[0], R0Y, C_STEEL,
                    ["IDLE", "wait start", "latch addr"]),
        "WR_ADDR": (CX[1], R0Y, C_NAVY,
                    ["WR_ADDR", "awvalid=1", "wvalid=1"]),
        "WR_RESP": (CX[2], R0Y, C_BLUE,
                    ["WR_RESP", "bready=1", "wait bvalid"]),
        "RD_ADDR": (CX[2], R1Y, C_LBLUE,
                    ["RD_ADDR", "arvalid=1", "wait arready"]),
        "RD_DATA": (CX[1], R1Y, C_TEAL,
                    ["RD_DATA", "rready=1", "wait rvalid"]),
        "DONE":    (CX[0], R1Y, C_GREEN,
                    ["DONE", "done=1 (1 cyc)", "assert error"]),
    }

    sbox = {}
    for name, (sx, sy, fill, lines) in STATES.items():
        sbox[name] = add_box(sl, sx, sy, BW, BH, lines, fill=fill, size=9.5)

    def cx(name): return STATES[name][0] + BW / 2
    def cy(name): return STATES[name][1] + BH / 2

    # ── Transition arrows ─────────────────────────────────────────────
    GAP = 0.12   # spacing between box edge and arrow

    # IDLE → WR_ADDR (right, row 0)
    add_arrow(sl, STATES["IDLE"][0]+BW+GAP, cy("IDLE"),
              STATES["WR_ADDR"][0]-GAP, cy("WR_ADDR"), color=C_ORANGE, width=2)
    add_txt(sl, "start=1", cx("IDLE")+0.72, cy("IDLE")-0.42, 1.1, 0.3,
            size=9, color=C_DGRAY, align=PP_ALIGN.CENTER)

    # WR_ADDR → WR_RESP (right, row 0)
    add_arrow(sl, STATES["WR_ADDR"][0]+BW+GAP, cy("WR_ADDR"),
              STATES["WR_RESP"][0]-GAP, cy("WR_RESP"), color=C_ORANGE, width=2)
    add_txt(sl, "awready\n&&wready", cx("WR_ADDR")+0.72, cy("WR_ADDR")-0.48, 1.1, 0.4,
            size=8.5, color=C_DGRAY, align=PP_ALIGN.CENTER)

    # WR_RESP → RD_ADDR (down, col 2)
    add_arrow(sl, cx("WR_RESP"), STATES["WR_RESP"][1]+BH+GAP,
              cx("RD_ADDR"), STATES["RD_ADDR"][1]-GAP, color=C_ORANGE, width=2)
    add_txt(sl, "bvalid", cx("WR_RESP")+0.12, (R0Y+R1Y)/2+0.35, 0.9, 0.3,
            size=9, color=C_DGRAY, align=PP_ALIGN.LEFT)

    # RD_ADDR → RD_DATA (left, row 1)
    add_arrow(sl, STATES["RD_ADDR"][0]-GAP, cy("RD_ADDR"),
              STATES["RD_DATA"][0]+BW+GAP, cy("RD_DATA"), color=C_ORANGE, width=2)
    add_txt(sl, "arready", cx("RD_DATA")+0.72, cy("RD_DATA")+0.32, 1.1, 0.3,
            size=9, color=C_DGRAY, align=PP_ALIGN.CENTER)

    # RD_DATA → DONE (left, row 1)
    add_arrow(sl, STATES["RD_DATA"][0]-GAP, cy("RD_DATA"),
              STATES["DONE"][0]+BW+GAP, cy("DONE"), color=C_ORANGE, width=2)
    add_txt(sl, "rvalid", cx("DONE")+0.72, cy("DONE")+0.32, 0.9, 0.3,
            size=9, color=C_DGRAY, align=PP_ALIGN.CENTER)

    # DONE → IDLE (up, col 0)
    add_arrow(sl, cx("DONE"), STATES["DONE"][1]-GAP,
              cx("IDLE"), STATES["IDLE"][1]+BH+GAP, color=C_GREEN, width=2)
    add_txt(sl, "unconditional", STATES["IDLE"][0]-1.55, (R0Y+R1Y)/2+0.25, 1.4, 0.3,
            size=9, bold=True, color=C_GREEN, align=PP_ALIGN.RIGHT)

    # ── Error logic annotation ────────────────────────────────────────
    add_rect(sl, 0.3, 5.0, 6.5, 0.75, fill=C_RED, rounded=True,
             lc=C_WHITE, lw=0)
    add_txt(sl, "error  =  write_err (bresp≠OKAY)   |   read_err (rresp≠OKAY)   |   data_err (rdata≠write_data)",
            0.4, 5.1, 6.3, 0.5,
            size=9.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ── Right: AXI output signal table ───────────────────────────────
    add_txt(sl, "AXI Output Signals by State", 6.8, 1.15, 6.3, 0.35,
            size=13, bold=True, color=C_NAVY)

    add_table(sl, 6.8, 1.55, 6.2, 3.65,
              ["State", "Driven Signals"],
              [["IDLE",    "all valids = 0"],
               ["WR_ADDR", "awvalid=1, wvalid=1, awaddr=addr_lat,\nwdata=wdata_lat, wstrb=0xF, wlast=1"],
               ["WR_RESP", "bready = 1"],
               ["RD_ADDR", "arvalid = 1,  araddr = addr_lat"],
               ["RD_DATA", "rready = 1"],
               ["DONE",    "done_r=1,  error_r = combined flags"]],
              hdr_fill=C_NAVY,
              alt_fill=RGBColor(0xE8, 0xF4, 0xF8))

    add_txt(sl, "All outputs are combinatorial from FSM state  ·  "
            "Data latches: addr_lat, wdata_lat, rdata_lat",
            6.8, 5.35, 6.2, 0.5,
            size=9, italic=True, color=C_DGRAY)

    add_rect(sl, 0, 6.75, 13.33, 0.75, fill=C_NAVY)
    add_txt(sl, "14 / 14 checks PASSED  ·  done pulses 1 cycle  ·  Simulation time 445 ns",
            0.3, 6.82, 12.7, 0.5,
            size=12, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)


slide_master_fsm()


# =============================================================================
# Slide 7 – Simple AXI4 Master: AXI Handshake Waveform & Test Results
# =============================================================================
def slide_master_timing():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, C_LGRAY)
    add_header(sl, "Simple AXI4 Master — AXI Handshake Waveform & Test Results")

    # ── Waveform section ──────────────────────────────────────────────
    add_txt(sl, "Single Write→Read Transaction  "
            "(FSM state shown in orange  ·  back-to-back handshakes, 6 active cycles)",
            0.3, 1.15, 12.73, 0.25,
            size=11, bold=True, color=C_NAVY)

    # Cycle index:  0       1        2        3        4        5     6     7   8  9
    state_labels = ["IDLE", "WR_ADDR","WR_RESP","RD_ADDR","RD_DATA","DONE","IDLE","·","·","·"]

    # Values verified against logs/simple_axi_master_tb.vcd (timescale 1 ps, 10 ns clk)
    # awready/wready/arready are PRE-ASSERTED by slave (HIGH when idle, dip LOW when busy)
    # start and awvalid overlap in the SAME cycle (both driven at/after posedge where
    # FSM latches start and transitions to WR_ADDR)
    SIG = [
        # (label,    [idx 0..9 values],               color   )
        ("start",   [0, 1, 0, 0, 0, 0, 0, 0, 0, 0],  C_ORANGE),
        ("awvalid", [0, 1, 0, 0, 0, 0, 0, 0, 0, 0],  C_TEAL),
        ("awready", [1, 1, 0, 1, 1, 1, 1, 1, 1, 1],  C_GREEN),
        ("wvalid",  [0, 1, 0, 0, 0, 0, 0, 0, 0, 0],  C_BLUE),
        ("wready",  [1, 1, 0, 1, 1, 1, 1, 1, 1, 1],  C_GREEN),
        ("bvalid",  [0, 0, 1, 0, 0, 0, 0, 0, 0, 0],  C_LBLUE),
        ("bready",  [0, 0, 1, 0, 0, 0, 0, 0, 0, 0],  C_ORANGE),
        ("arvalid", [0, 0, 0, 1, 0, 0, 0, 0, 0, 0],  C_TEAL),
        ("arready", [1, 1, 1, 1, 0, 1, 1, 1, 1, 1],  C_GREEN),
        ("rvalid",  [0, 0, 0, 0, 1, 0, 0, 0, 0, 0],  C_LBLUE),
        ("rready",  [0, 0, 0, 0, 1, 0, 0, 0, 0, 0],  C_ORANGE),
        ("done",    [0, 0, 0, 0, 0, 0, 1, 0, 0, 0],  C_GREEN),
    ]

    # Waveform: y=1.42 to y=5.30  →  height=3.88
    add_waveform(sl, 0.3, 1.42, 12.73, 3.88, SIG, state_labels,
                 label_w=1.1, max_font=8.5)

    # ── Test results table ────────────────────────────────────────────
    add_txt(sl, "Test Results  (14 Checks)",
            0.3, 5.40, 12.5, 0.27,
            size=12, bold=True, color=C_NAVY)

    add_table(sl, 0.3, 5.69, 12.5, 1.65,
              ["Test", "Address", "Write Data", "done", "error", "read_data", "Result"],
              [["T1: Single W/R",  "0x0000_0000", "0xDEAD_BEEF", "1", "0", "0xDEAD_BEEF", "PASS"],
               ["T2a",             "0x0000_0004", "0xCAFE_F00D", "1", "0", "0xCAFE_F00D", "PASS"],
               ["T2b",             "0x0000_0008", "0xA5A5_A5A5", "1", "0", "0xA5A5_A5A5", "PASS"],
               ["T2c",             "0x0000_000C", "0x0000_0001", "1", "0", "0x0000_0001", "PASS"],
               ["T3: OOB 0x1000",  "0x0000_1000", "0x1234_5678", "1", "1", "SLVERR",      "PASS"]],
              hdr_fill=C_NAVY,
              alt_fill=RGBColor(0xE8, 0xF4, 0xF8))


slide_master_timing()


# =============================================================================
# Slide 8 – Summary
# =============================================================================
def slide_summary():
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, C_NAVY)

    add_rect(sl, 0, 0, 13.33, 1.05, fill=C_TEAL)
    add_txt(sl, "Summary — Overall Verification Status",
            0.35, 0.12, 12.6, 0.82,
            size=30, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # Three component cards
    for (title, file_, count, bullets), x in zip(
        [
            ("AXI-Lite\nRegister File",
             "rtl/axi_lite_regfile.v",
             "133 / 133",
             ["4 × 32-bit registers",
              "SLVERR on addr[31:4]≠0",
              "3 write FSM paths",
              "Backpressure coverage"]),
            ("AXI4\nMemory Model",
             "tb/axi_mem_model.sv",
             "30 / 30",
             ["1 KB parameterised RAM",
              "Byte-WSTRB  do_write()",
              "SLVERR for OOB addresses",
              "Init to zero at sim start"]),
            ("Simple\nAXI4 Master",
             "rtl/simple_axi_master.v",
             "14 / 14",
             ["6-state FSM",
              "Write → Readback sequence",
              "Error: SLVERR / data mismatch",
              "Combinatorial AXI drivers"]),
        ],
        [0.45, 4.6, 8.75]):

        CW = 4.0
        # Card background
        add_rect(sl, x, 1.15, CW, 5.7,
                 fill=RGBColor(0x1E, 0x3A, 0x5A),
                 lc=C_TEAL, lw=1.0, rounded=True)

        add_txt(sl, title, x + 0.1, 1.25, CW - 0.2, 0.65,
                size=17, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)

        add_txt(sl, file_, x + 0.1, 1.88, CW - 0.2, 0.32,
                size=9, italic=True,
                color=RGBColor(0x99, 0xCC, 0xDD), align=PP_ALIGN.CENTER)

        # Check count badge
        add_rect(sl, x + 0.4, 2.3, CW - 0.8, 0.85,
                 fill=C_GREEN, rounded=True)
        add_txt(sl, count, x + 0.4, 2.33, CW - 0.8, 0.58,
                size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_txt(sl, "checks passed", x + 0.4, 2.88, CW - 0.8, 0.28,
                size=10, color=C_LGRAY, align=PP_ALIGN.CENTER)

        # Bullet points
        for i, b in enumerate(bullets):
            add_txt(sl, "▸  " + b,
                    x + 0.25, 3.35 + i * 0.55, CW - 0.35, 0.45,
                    size=12, color=C_LGRAY)

    # ── Bottom total bar ──────────────────────────────────────────────
    add_rect(sl, 0.45, 6.95, 12.43, 0.45, fill=C_ORANGE, rounded=True)
    add_txt(sl,
            "Total:  177 / 177 checks passed  ·  0 failures  ·  "
            "3 CI scripts  ·  All exit code 0",
            0.45, 6.97, 12.43, 0.38,
            size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


slide_summary()


# =============================================================================
# Save
# =============================================================================
os.makedirs("slides", exist_ok=True)
OUT = "slides/project_overview.pptx"
prs.save(OUT)
print(f"Saved  : {OUT}")
print(f"Slides : {len(prs.slides)}")
print("All elements are native PowerPoint shapes — fully editable.")
